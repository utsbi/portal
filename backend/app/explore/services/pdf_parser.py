from typing import List, Dict, Any, Optional
from pypdf import PdfReader
from fastapi import HTTPException
import io
import os
import zipfile


# Plain-text / markup formats whose bytes can be decoded directly as UTF-8.
_PLAINTEXT_EXTENSIONS = {".txt", ".md"}

# ---------------------------------------------------------------------------
# Zip-bomb guard — Office formats (docx/pptx/xlsx) are zip archives.
# We inspect the Central Directory before handing bytes to the Office
# parsers, so a malicious archive cannot expand to gigabytes in memory.
# ---------------------------------------------------------------------------

# Maximum total *declared* uncompressed size across all members (200 MB).
# Real Office documents rarely exceed 50 MB decompressed; 200 MB is generous.
_ZIP_MAX_DECOMPRESSED_BYTES: int = 200 * 1024 * 1024

# Maximum number of members in the archive.  A normal .docx has ~20-50
# entries; 2 000 is very generous and still blocks archive-flooding attacks.
_ZIP_MAX_MEMBER_COUNT: int = 2_000

# Maximum compression ratio: total_uncompressed / on-wire archive size.
# Legitimate text-heavy documents rarely compress beyond 10:1; 100 is a
# conservative ceiling that catches classic nested-deflate zip bombs.
_ZIP_MAX_RATIO: int = 100

# Binary document formats we have a dedicated extractor for.
_DOCX_EXTENSIONS = {".docx"}
_PPTX_EXTENSIONS = {".pptx"}
_XLSX_EXTENSIONS = {".xlsx"}
_PDF_EXTENSIONS = {".pdf"}

# Every extension the dispatcher can turn into plain text. Anything else
# (images, archives, audio, ...) is "not indexable" — not an error.
SUPPORTED_EXTENSIONS = (
    _PLAINTEXT_EXTENSIONS
    | _DOCX_EXTENSIONS
    | _PPTX_EXTENSIONS
    | _XLSX_EXTENSIONS
    | _PDF_EXTENSIONS
)


def _check_zip_bomb(file_bytes: bytes) -> None:
    """Inspect a zip archive's Central Directory for zip-bomb signatures.

    Reads only metadata (no decompression) and raises ``HTTPException(400)``
    if any of the three limits is exceeded:

    * Total *declared* uncompressed size across all members > ``_ZIP_MAX_DECOMPRESSED_BYTES``
    * Member count > ``_ZIP_MAX_MEMBER_COUNT``
    * Compression ratio (total_uncompressed / archive_size) > ``_ZIP_MAX_RATIO``

    The Central Directory can lie about sizes, but inflated declarations are
    themselves the attack vector (they consume memory in the parsers that trust
    them), so checking declared sizes is the correct defence.

    If the bytes are not a valid zip at all we return silently and let the
    downstream Office parser produce its own error.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            members = zf.infolist()
    except zipfile.BadZipFile:
        return

    if len(members) > _ZIP_MAX_MEMBER_COUNT:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Archive rejected: member count {len(members)} exceeds "
                f"limit of {_ZIP_MAX_MEMBER_COUNT}"
            ),
        )

    total_uncompressed = sum(m.file_size for m in members)
    if total_uncompressed > _ZIP_MAX_DECOMPRESSED_BYTES:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Archive rejected: declared decompressed size "
                f"{total_uncompressed} bytes exceeds limit of "
                f"{_ZIP_MAX_DECOMPRESSED_BYTES} bytes"
            ),
        )

    archive_size = len(file_bytes)
    if archive_size > 0 and total_uncompressed > archive_size * _ZIP_MAX_RATIO:
        ratio = total_uncompressed / archive_size
        raise HTTPException(
            status_code=400,
            detail=(
                f"Archive rejected: compression ratio {ratio:.1f}x exceeds "
                f"limit of {_ZIP_MAX_RATIO}x"
            ),
        )


def _file_extension(filename: str) -> str:
    """Lower-cased extension (including the dot) for ``filename``, or ``""``."""
    return os.path.splitext(filename)[1].lower()


def is_extractable(filename: str) -> bool:
    """True if a file with this name has a text extractor (by extension)."""
    return _file_extension(filename) in SUPPORTED_EXTENSIONS


def _extract_plaintext(file_bytes: bytes) -> str:
    """Decode raw text/markdown bytes as UTF-8 (replacing undecodable bytes)."""
    return file_bytes.decode("utf-8", errors="replace")


def _extract_pdf(file_bytes: bytes) -> str:
    """Concatenate the text of every PDF page into a single string."""
    reader = PdfReader(io.BytesIO(file_bytes))
    parts: List[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(file_bytes: bytes) -> str:
    """Extract paragraph text from a .docx (Word) document."""
    _check_zip_bomb(file_bytes)
    from docx import Document

    document = Document(io.BytesIO(file_bytes))
    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from every shape across all slides of a .pptx deck."""
    _check_zip_bomb(file_bytes)
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(file_bytes))
    parts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = getattr(shape, "text_frame", None)
                if text_frame is None:
                    continue
                text = text_frame.text
                if text and text.strip():
                    parts.append(text)
    return "\n".join(parts)


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract cell values from every sheet of a .xlsx workbook as TSV rows."""
    _check_zip_bomb(file_bytes)
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    workbook.close()
    return "\n".join(parts)


def extract_text(file_bytes: bytes, filename: str) -> Optional[str]:
    """Extract plain text from ``file_bytes`` based on ``filename``'s extension.

    Dispatches to the right extractor for PDF, .txt, .md, .docx, .pptx and
    .xlsx. Returns ``None`` for any unsupported/binary type (images, archives,
    ...) — those are "not indexable", which the caller treats as a normal
    outcome rather than an error.
    """
    ext = _file_extension(filename)
    if ext in _PLAINTEXT_EXTENSIONS:
        return _extract_plaintext(file_bytes)
    if ext in _PDF_EXTENSIONS:
        return _extract_pdf(file_bytes)
    if ext in _DOCX_EXTENSIONS:
        return _extract_docx(file_bytes)
    if ext in _PPTX_EXTENSIONS:
        return _extract_pptx(file_bytes)
    if ext in _XLSX_EXTENSIONS:
        return _extract_xlsx(file_bytes)
    return None


class PDFParser:
    """Service for extracting text and metadata from PDF files."""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_text_with_metadata(
        self, file_bytes: bytes, filename: str
    ) -> List[Dict[str, Any]]:
        """Extract text from a PDF file with page-level metadata."""
        pages_data = []

        try:
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            total_pages = len(pdf_reader.pages)

            for page_num, page in enumerate(pdf_reader.pages, start=1):
                text = page.extract_text()

                if text and text.strip():
                    pages_data.append(
                        {
                            "content": self._clean_text(text),
                            "metadata": {
                                "filename": filename,
                                "page_number": page_num,
                                "total_pages": total_pages,
                                "file_type": "pdf",
                            },
                        }
                    )

        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")

        return pages_data

    def chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks for embedding."""
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        chunks = splitter.split_text(text)
        return [chunk.strip() for chunk in chunks if chunk.strip()]

    def chunk_pages(self, pages_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Chunk extracted page data into smaller pieces for embedding."""
        chunked_data = []

        for page_data in pages_data:
            text = page_data["content"]
            base_metadata = page_data["metadata"]

            chunks = self.chunk_text(text)

            for chunk_idx, chunk in enumerate(chunks):
                chunked_data.append(
                    {
                        "content": chunk,
                        "metadata": {
                            **base_metadata,
                            "chunk_index": chunk_idx,
                            "total_chunks_in_page": len(chunks),
                        },
                    }
                )

        return chunked_data

    def _clean_text(self, text: str) -> str:
        """Clean extracted text by normalizing whitespace."""
        lines = text.split("\n")
        cleaned_lines = []

        for line in lines:
            cleaned_line = " ".join(line.split())
            if cleaned_line:
                cleaned_lines.append(cleaned_line)

        return "\n".join(cleaned_lines)
